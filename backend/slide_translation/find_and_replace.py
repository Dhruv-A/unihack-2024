import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from .translator import Translator
from .img_to_txt import ImageToText

n=0
class FindAndReplace:
    def __init__(self):
        self.translator = Translator()

    def annotate_image_in_slide(self, image, slide, language, file_name):
        global n
        original_image = image
        img_bytes = original_image.blob
        filename = 'image{:03d}.{}'.format(n, image.ext)
        filepath = f'slide_assets/{file_name}'
        n += 1

        if not os.path.exists(filepath):
            os.makedirs(filepath)

        filepath = os.path.join(f'slide_assets/{file_name}', filename)

        with open(filepath, 'wb') as file:
            file.write(img_bytes)
        
        img2text = ImageToText(filepath, os.getenv('API_NINJA_API_KEY'))
        imgtext = img2text.get_concat_string()
        translated_text = self.translator.translate(imgtext, language)

        # Add image text to presenter notes
        notes_slide = slide.notes_slide

        text_frame = notes_slide.notes_text_frame
        text_frame.text = text_frame.text + '\n' + f'Image translation: ${translated_text}'

    def replace_text_in_slide(self, language, slide, file_name):
        """
        Replaces text in a slide with its corresponding translation.
        'replacements' is a dictionary where keys are original texts and values are their translations.
        """
        for shape in slide.shapes:
            # Translate well annotated slide parts
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for s in shape.shapes:
                    self.annotate_image_in_slide(shape.image, slide, language, file_name)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                self.annotate_image_in_slide(shape.image, slide, language, file_name)    
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    translated_text = self.translator.translate(original_text, language)
                    run.text = translated_text

    def replace_text_in_presentation(self, file_name, language, presentation_path):
        """
        Loads a presentation, replaces text on each slide, and saves the presentation.
        """
        prs = Presentation(presentation_path)
        for slide in prs.slides:
            self.replace_text_in_slide(language, slide, file_name)

        if not os.path.exists('translated_presentations'):
            os.makedirs('translated_presentations')
        prs.save(f'translated_presentations/{file_name}_translated.pptx')
